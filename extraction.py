import json
import os
import tempfile

import config

_whisper = None


def extraire_texte(chemin):
    extension = os.path.splitext(chemin)[1].lower()
    lecteur = _REGISTRE.get(extension)
    if lecteur is None:
        if extension in config.EXTENSIONS_IMAGE:
            lecteur = _lire_image
        elif extension in config.EXTENSIONS_AUDIO:
            lecteur = _lire_audio
        elif extension in config.EXTENSIONS_VIDEO:
            lecteur = _lire_video
        else:
            return ""
    try:
        return (lecteur(chemin) or "").strip()
    except Exception as erreur:
        print(f"  ! Lecture impossible ({chemin}) : {erreur}")
        return ""


def _lire_pdf(chemin):
    from pypdf import PdfReader

    lecteur = PdfReader(chemin)
    texte = "\n".join((page.extract_text() or "") for page in lecteur.pages)
    if texte.strip():
        return texte
    return _ocr_pdf(chemin)


def _ocr_pdf(chemin):
    try:
        from pdf2image import convert_from_path
    except ImportError:
        return ""
    morceaux = []
    for image in convert_from_path(chemin):
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as fichier:
            cible = fichier.name
        image.save(cible)
        morceaux.append(_ocr_image(cible))
        os.remove(cible)
    return "\n".join(m for m in morceaux if m)


def _lire_docx(chemin):
    import docx

    document = docx.Document(chemin)
    paragraphes = [p.text for p in document.paragraphs]
    for table in document.tables:
        for ligne in table.rows:
            paragraphes.append(" | ".join(cellule.text for cellule in ligne.cells))
    return "\n".join(paragraphes)


def _lire_pptx(chemin):
    from pptx import Presentation

    presentation = Presentation(chemin)
    morceaux = []
    for diapositive in presentation.slides:
        for forme in diapositive.shapes:
            if forme.has_text_frame:
                morceaux.append(forme.text_frame.text)
    return "\n".join(morceaux)


def _lire_xlsx(chemin):
    from openpyxl import load_workbook

    classeur = load_workbook(chemin, read_only=True, data_only=True)
    morceaux = []
    for feuille in classeur.worksheets:
        morceaux.append(f"# {feuille.title}")
        for ligne in feuille.iter_rows(values_only=True):
            valeurs = [str(v) for v in ligne if v is not None]
            if valeurs:
                morceaux.append(" | ".join(valeurs))
    return "\n".join(morceaux)


def _lire_csv(chemin):
    import csv

    with open(chemin, newline="", encoding="utf-8", errors="ignore") as fichier:
        lecteur = csv.reader(fichier)
        return "\n".join(" | ".join(ligne) for ligne in lecteur)


def _lire_html(chemin):
    from bs4 import BeautifulSoup

    with open(chemin, encoding="utf-8", errors="ignore") as fichier:
        soupe = BeautifulSoup(fichier.read(), "html.parser")
    for balise in soupe(["script", "style"]):
        balise.decompose()
    return soupe.get_text(separator="\n")


def _lire_json(chemin):
    with open(chemin, encoding="utf-8", errors="ignore") as fichier:
        donnees = json.load(fichier)
    return json.dumps(donnees, ensure_ascii=False, indent=2)


def _lire_rtf(chemin):
    from striprtf.striprtf import rtf_to_text

    with open(chemin, encoding="utf-8", errors="ignore") as fichier:
        return rtf_to_text(fichier.read())


def _lire_texte(chemin):
    with open(chemin, encoding="utf-8", errors="ignore") as fichier:
        return fichier.read()


def _lire_image(chemin):
    parties = []
    texte_ocr = _ocr_image(chemin)
    if texte_ocr:
        parties.append("Texte détecté dans l'image :\n" + texte_ocr)
    description = _decrire_image(chemin)
    if description:
        parties.append("Description de l'image :\n" + description)
    return "\n\n".join(parties)


def _ocr_image(chemin):
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""
    try:
        return pytesseract.image_to_string(Image.open(chemin), lang="fra+eng").strip()
    except Exception:
        try:
            return pytesseract.image_to_string(Image.open(chemin)).strip()
        except Exception:
            return ""


def _decrire_image(chemin):
    try:
        import ollama
    except ImportError:
        return ""
    try:
        reponse = ollama.chat(
            model=config.VISION_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": "Décris en détail et en français le contenu de cette image, y compris tout texte visible.",
                    "images": [chemin],
                }
            ],
        )
        return reponse["message"]["content"].strip()
    except Exception:
        return ""


def _lire_audio(chemin):
    return _transcrire(chemin)


def _lire_video(chemin):
    parties = []
    transcription = _transcrire(chemin)
    if transcription:
        parties.append("Transcription audio :\n" + transcription)
    descriptions = []
    for image in _extraire_images(chemin):
        description = _decrire_image(image)
        if description:
            descriptions.append(description)
        os.remove(image)
    if descriptions:
        parties.append("Scènes de la vidéo :\n" + "\n\n".join(descriptions))
    return "\n\n".join(parties)


def _charger_whisper():
    global _whisper
    if _whisper is None:
        from faster_whisper import WhisperModel

        _whisper = WhisperModel(
            config.WHISPER_MODEL,
            device=config.WHISPER_DEVICE,
            compute_type=config.WHISPER_COMPUTE,
        )
    return _whisper


def _transcrire(chemin):
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        return ""
    try:
        modele = _charger_whisper()
        segments, _ = modele.transcribe(chemin)
        return " ".join(segment.text.strip() for segment in segments).strip()
    except Exception as erreur:
        print(f"  ! Transcription impossible ({chemin}) : {erreur}")
        return ""


def _extraire_images(chemin):
    nombre = config.NB_IMAGES_VIDEO
    if nombre <= 0:
        return []
    try:
        import av
    except ImportError:
        return []
    images = []
    try:
        conteneur = av.open(chemin)
        flux = conteneur.streams.video[0]
        flux.thread_type = "AUTO"
        duree = 0.0
        if flux.duration and flux.time_base:
            duree = float(flux.duration * flux.time_base)
        elif conteneur.duration:
            duree = conteneur.duration / 1_000_000.0
        for i in range(nombre):
            instant = duree * (i + 1) / (nombre + 1) if duree > 0 else 0.0
            if instant > 0 and flux.time_base:
                try:
                    conteneur.seek(int(instant / flux.time_base), stream=flux)
                except Exception:
                    pass
            image = _premiere_image(conteneur, flux)
            if image is not None:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as fichier:
                    cible = fichier.name
                image.to_image().save(cible)
                images.append(cible)
        conteneur.close()
    except Exception as erreur:
        print(f"  ! Extraction d'images impossible ({chemin}) : {erreur}")
    return images


def _premiere_image(conteneur, flux):
    for trame in conteneur.decode(flux):
        return trame
    return None


_REGISTRE = {
    ".pdf": _lire_pdf,
    ".docx": _lire_docx,
    ".pptx": _lire_pptx,
    ".xlsx": _lire_xlsx,
    ".xlsm": _lire_xlsx,
    ".csv": _lire_csv,
    ".html": _lire_html,
    ".htm": _lire_html,
    ".json": _lire_json,
    ".rtf": _lire_rtf,
    ".txt": _lire_texte,
    ".md": _lire_texte,
    ".markdown": _lire_texte,
    ".log": _lire_texte,
    ".xml": _lire_texte,
    ".yaml": _lire_texte,
    ".yml": _lire_texte,
}
